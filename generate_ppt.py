"""
SentriX — Hackathon 2026 Presentation Generator
Clones the template PPTX and updates content while strictly preserving design.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import copy
import os

TEMPLATE = os.path.join(os.path.dirname(__file__),
    "Idea Submission Solution  - Template_Hachsplosion 2026.pptx (1).pptx")
OUTPUT = os.path.join(os.path.dirname(__file__),
    "SentriX - Hackathon 2026 Submission.pptx")

# Enterprise Core palette
GREEN  = RGBColor(0x86, 0xBC, 0x25)
AMBER  = RGBColor(0xF4, 0xB4, 0x3C)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLACK  = RGBColor(0x00, 0x00, 0x00)
BLUE   = RGBColor(0x3F, 0x6C, 0xE4)
DGREEN = RGBColor(0x26, 0x89, 0x0D)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def iter_all_shapes(slide):
    """Yield every shape on the slide, recursively entering groups."""
    for shape in slide.shapes:
        yield shape
        if shape.shape_type == 6:  # GROUP
            yield from _iter_group(shape)

def _iter_group(group):
    for shape in group.shapes:
        yield shape
        if shape.shape_type == 6:
            yield from _iter_group(shape)

def replace_in_runs(shape, old, new):
    """Replace text in every run of every paragraph, preserving formatting."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)

def set_textbox_text_preserve_first_format(shape, new_text, font_size=None,
                                            bold=None, color=None, alignment=None):
    """Clear a text box and set simple single-run text, reusing first run's format."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # Capture first run's format
    ref_font = None
    for para in tf.paragraphs:
        if para.runs:
            ref_font = para.runs[0].font
            break

    # Clear all paragraphs except first
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    para = tf.paragraphs[0]
    # Clear existing runs
    for r in list(para.runs):
        para._p.remove(r._r)

    # Handle multi-line
    lines = new_text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            run = para.add_run()
        else:
            from pptx.oxml.ns import qn
            import lxml.etree as etree
            # Add new paragraph
            new_p = copy.deepcopy(para._p)
            # Remove runs from new_p
            for r in new_p.findall(qn('a:r')):
                new_p.remove(r)
            tf._txBody.append(new_p)
            from pptx.text.text import _Paragraph
            new_para = _Paragraph(new_p, tf)
            run = new_para.add_run()

        run.text = line
        # Apply format
        if ref_font:
            if font_size:
                run.font.size = Pt(font_size)
            elif ref_font.size:
                run.font.size = ref_font.size

            if bold is not None:
                run.font.bold = bold
            elif ref_font.bold:
                run.font.bold = ref_font.bold

            if color:
                run.font.color.rgb = color
            elif ref_font.color and ref_font.color.rgb:
                run.font.color.rgb = ref_font.color.rgb

            if ref_font.name:
                run.font.name = ref_font.name
        else:
            if font_size:
                run.font.size = Pt(font_size)
            if bold is not None:
                run.font.bold = bold
            if color:
                run.font.color.rgb = color

    if alignment:
        para.alignment = alignment


def set_multiformat_text(shape, segments):
    """
    Set text with multiple formatting segments.
    segments: list of (text, color, bold) tuples
    Preserves font name and size from the first existing run.
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # Capture reference format
    ref_size = None
    ref_name = None
    for para in tf.paragraphs:
        if para.runs:
            ref_size = para.runs[0].font.size
            ref_name = para.runs[0].font.name
            break

    # Clear all but first paragraph
    while len(tf.paragraphs) > 1:
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)

    para = tf.paragraphs[0]
    for r in list(para.runs):
        para._p.remove(r._r)

    for text, color, bold in segments:
        if '\n' in text:
            parts = text.split('\n')
            for j, part in enumerate(parts):
                if j == 0:
                    run = para.add_run()
                    run.text = part
                else:
                    from pptx.oxml.ns import qn
                    new_p = copy.deepcopy(para._p)
                    for r in new_p.findall(qn('a:r')):
                        new_p.remove(r)
                    tf._txBody.append(new_p)
                    from pptx.text.text import _Paragraph
                    para = _Paragraph(new_p, tf)
                    run = para.add_run()
                    run.text = part

                run.font.bold = bold
                if color:
                    run.font.color.rgb = color
                if ref_size:
                    run.font.size = ref_size
                if ref_name:
                    run.font.name = ref_name
        else:
            run = para.add_run()
            run.text = text
            run.font.bold = bold
            if color:
                run.font.color.rgb = color
            if ref_size:
                run.font.size = ref_size
            if ref_name:
                run.font.name = ref_name


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation(TEMPLATE)

    # ===== GLOBAL REPLACEMENTS across all slides =====
    for slide in prs.slides:
        for shape in iter_all_shapes(slide):
            replace_in_runs(shape, "ComplianceGPT", "SentriX")
            replace_in_runs(shape, "complianceGPT", "SentriX")
            replace_in_runs(shape, "complaince_gpt", "SentriX-platform")

    # ===== SLIDE-BY-SLIDE CONTENT UPDATES =====

    # --- Slide 2 (index 1): Team Intro ---
    slide2 = prs.slides[1]
    for shape in iter_all_shapes(slide2):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if t == "Text boxes with arrows":
            set_textbox_text_preserve_first_format(shape, "Team Introduction")
        if t == "Impacters":
            set_textbox_text_preserve_first_format(shape, "Team Impacters")
        if "Forensic and Financial Crime" in t:
            set_textbox_text_preserve_first_format(
                shape,
                "Forensic & Financial Crime",
                color=WHITE)

    # --- Slide 3 (index 2): Problem Definition ---
    slide3 = prs.slides[2]
    for shape in iter_all_shapes(slide3):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        # Main problem text
        if "The Primary issue" in t:
            set_multiformat_text(shape, [
                ("The ", WHITE, False),
                ("$2.5 Trillion", GREEN, True),
                (" global compliance market is drowning in ", WHITE, False),
                ("manual processes", GREEN, True),
                (". Auditors spend ", WHITE, False),
                ("4-12 weeks", GREEN, True),
                (" per assessment, manually evaluating ", WHITE, False),
                ("Excel spreadsheets ", GREEN, True),
                ("of ISO clauses against scattered policy documents. The result? ", WHITE, False),
                ("Inconsistent scores", GREEN, True),
                (", ", WHITE, False),
                ("missed gaps", GREEN, True),
                (", ", WHITE, False),
                ("$50K-$500K costs", GREEN, True),
                (" per engagement, and reports that are ", WHITE, False),
                ("outdated on delivery", GREEN, True),
                (". Organizations managing ", WHITE, False),
                ("4+ ISO standards", GREEN, True),
                (" simultaneously face compounded ", WHITE, False),
                ("delays", GREEN, True),
                (" and ", WHITE, False),
                ("duplicated effort", GREEN, True),
                (".", WHITE, False),
            ])

        if "Unstructured Policy Document" in t:
            set_multiformat_text(shape, [
                ("Unstructured Policy Documents", AMBER, False),
                (" — Governance frameworks are scattered across PDFs, Word docs, and spreadsheets with no standardized structure for AI analysis.", WHITE, False),
                ("\n\n", WHITE, False),
                ("Complex Multi-Standard ISO Mapping", AMBER, False),
                (" — ISO frameworks contain hundreds of clauses with subtle interdependencies that must be accurately mapped against evidence.", WHITE, False),
                ("\n\n", WHITE, False),
                ("Manual & Time-Consuming Assessments", AMBER, False),
                (" — Excel-based checklist reviews consume 4–12 weeks of senior consultant time at $300-$500/hr.", WHITE, False),
                ("\n\n", WHITE, False),
                ("Evidence Validation Gap", AMBER, False),
                (" — No existing tool verifies whether cited evidence actually proves compliance — auditors catch this, AI tools don't.", WHITE, False),
                ("\n\n", WHITE, False),
                ("Scalability Across Standards", AMBER, False),
                (" — Organizations duplicate effort when the same policy satisfies clauses across ISO 27001, 37001, 9001, and more.", WHITE, False),
            ])

        if "Cost Impact" in t:
            set_multiformat_text(shape, [
                ("Cost Impact", AMBER, True),
                (" — Traditional compliance consulting: ", WHITE, False),
                ("$50,000 – $500,000", GREEN, True),
                (" per assessment. SentriX reduces this by ", WHITE, False),
                ("90%+", GREEN, True),
                (".", WHITE, False),
            ])

        if "Time Impact" in t:
            set_multiformat_text(shape, [
                ("Time Impact", AMBER, True),
                (" — Manual assessments: ", WHITE, False),
                ("4–12 weeks", GREEN, True),
                (" with teams of auditors. SentriX delivers in ", WHITE, False),
                ("minutes", GREEN, True),
                (".", WHITE, False),
            ])

        if "Accuracy Issues" in t:
            set_multiformat_text(shape, [
                ("Accuracy Crisis", AMBER, True),
                (" — Two human auditors, same document = ", WHITE, False),
                ("two different scores", GREEN, True),
                (". SentriX's 3-Tier engine ensures ", WHITE, False),
                ("100% reproducibility", GREEN, True),
                (".", WHITE, False),
            ])

        if "Lack of Visibility" in t:
            set_multiformat_text(shape, [
                ("\n", WHITE, False),
                ("No Real-Time Visibility", AMBER, True),
                (" — Static PDF reports are obsolete on delivery, with ", WHITE, False),
                ("zero interactive drill-down", GREEN, True),
                (" into clause-level gaps.", WHITE, False),
            ])

    # --- Slide 4 (index 3): Solution Approach ---
    slide4 = prs.slides[3]
    for shape in iter_all_shapes(slide4):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        if "Project Title" in t:
            set_multiformat_text(shape, [
                ("Project Title", AMBER, True),
                (" — ", WHITE, False),
                ("SentriX", GREEN, True),
                (" (Agentic AI-Powered Multi-Standard Compliance Intelligence Platform built using ", WHITE, False),
                ("Gen", WHITE, True),
                ("W", GREEN, True),
                (".AI", WHITE, True),
                (")", WHITE, False),
            ])

        if "We are addressing" in t:
            set_multiformat_text(shape, [
                ("SentriX ", GREEN, True),
                ("transforms manual Excel-based ISO assessments into an ", WHITE, False),
                ("AI-powered enterprise platform", GREEN, False),
                (" that analyzes ", WHITE, False),
                ("uploaded policy documents", GREEN, False),
                (", runs a ", WHITE, False),
                ("7-agent orchestrated pipeline", GREEN, False),
                (", and delivers ", WHITE, False),
                ("instant compliance readiness scores", GREEN, False),
                (", ", WHITE, False),
                ("evidence validation", GREEN, False),
                (", ", WHITE, False),
                ("gap analysis", GREEN, False),
                (", and ", WHITE, False),
                ("auto-generated remediation policies", GREEN, False),
                (" — across ", WHITE, False),
                ("6 ISO standards simultaneously", GREEN, True),
                (" — all powered by ", WHITE, False),
                ("Enterprise Core's Gen", WHITE, True),
                ("W", GREEN, True),
                (".AI", WHITE, True),
                (".\n\n", WHITE, False),
                ("Our USP —", WHITE, True),
                ("\n\n", WHITE, False),
                ("3-Tier Hybrid Scoring Engine", AMBER, True),
                (" — ML Semantic + Groq AI + Keyword NLP ensures ", WHITE, False),
                ("reproducible, bias-free, audit-ready scores", GREEN, False),
                (" that degrade gracefully. No single point of failure.\n\n", WHITE, False),
                ("Evidence-Backed Compliance Verification", AMBER, True),
                (" — Industry-first Evidence Validation Agent that classifies evidence sufficiency (sufficient → partial → missing) and quality (direct → indirect → anecdotal).\n\n", WHITE, False),
                ("Real-Time AI Copilot", AMBER, True),
                (" — Assessment-aware conversational AI that explains scores, generates executive summaries, and provides remediation guidance in plain English.\n\n", WHITE, False),
                ("Cross-Standard Synergy Detection", AMBER, True),
                (" — Identifies where a single policy fix satisfies requirements across multiple ISO standards, eliminating duplicated effort.", WHITE, False),
            ])

    # --- Slide 5 (index 4): OmniAgent Modules ---
    slide5 = prs.slides[4]
    for shape in iter_all_shapes(slide5):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        if "It is time to unleash" in t:
            set_multiformat_text(shape, [
                ("Our architecture is purpose-built for ", AMBER, True),
                ("Gen", WHITE, True),
                ("W", GREEN, True),
                (".AI", WHITE, True),
                (". ", AMBER, True),
                ("Modules like ", AMBER, True),
                ("AppMaker", GREEN, True),
                (", ", AMBER, True),
                ("Realm AI", GREEN, True),
                (", ", AMBER, True),
                ("Agent Builder", GREEN, True),
                (" and ", AMBER, True),
                ("Playground", GREEN, True),
                (" power SentriX's automated AI-driven ISO compliance assessment pipeline from upload to remediation.", AMBER, True),
            ])

        if shape.name == "TextBox 16" and "We will build" in t:
            set_multiformat_text(shape, [
                ("SentriX's enterprise interface is built for AppMaker deployment.\n", WHITE, False),
                ("Rapid creation of document upload flows, compliance questionnaires, and multi-standard assessment configuration.\n", WHITE, False),
                ("Handles deployment and seamless integration of all backend AI services.", WHITE, False),
            ])

        if shape.name == "TextBox 18" and "Realm AI" in t:
            set_multiformat_text(shape, [
                ("Realm AI powers SentriX's connection to multiple LLMs (Groq, OpenAI) for compliance intelligence.\n", WHITE, False),
                ("Executes document interpretation, ISO clause mapping, KB retrieval, and compliance reasoning.\n", WHITE, False),
                ("Powers the critical 3-Tier Hybrid Scoring and Evidence Validation engines.", WHITE, False),
            ])

        if shape.name == "TextBox 20" and "Agent Builder" in t:
            set_multiformat_text(shape, [
                ("The soul of SentriX — Agent Builder orchestrates our 7-agent compliance pipeline.\n", WHITE, False),
                ("What the Agent Builder powers:\n\n", WHITE, False),
                (" Document Parsing Agent         Evidence Validation Agent\n", GREEN, False),
                (" Clause Mapping Agent            Compliance Scoring Agent\n", GREEN, False),
                (" Gap Detection Agent              Remediation Planning Agent\n", GREEN, False),
                (" Policy Generation Agent        AI Compliance Copilot\n", GREEN, False),
                ("Enables full automation of the end-to-end compliance assessment workflow.", WHITE, False),
            ])

        if shape.name == "TextBox 22" and "Playground" in t:
            set_multiformat_text(shape, [
                ("Playground powers SentriX's Interactive Compliance Analytics Dashboard.\n", WHITE, False),
                ("Visualizes our novel Organizational Risk Heatmap, Clause Heat Maps, and Gap Priority Matrices.\n", WHITE, False),
                ("Enables stakeholders to interactively explore compliance metrics with drill-down into clause-level detail.\n", WHITE, False),
                ("Not just visualization — enables scenario analysis and compliance forecasting.", WHITE, False),
            ])

    # --- Slide 6 (index 5): Step-by-Step Workflow ---
    slide6 = prs.slides[5]
    for shape in iter_all_shapes(slide6):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if "Step by Step" in t:
            set_textbox_text_preserve_first_format(
                shape,
                "Step-by-Step Workflow — How SentriX Transforms Compliance in Minutes",
                color=AMBER, bold=True)

    # --- Slide 7 (index 6): ER Model ---
    slide7 = prs.slides[6]
    for shape in iter_all_shapes(slide7):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if "ER Model" in t:
            set_textbox_text_preserve_first_format(
                shape, "ER Model for SentriX", color=AMBER, bold=True)

    # --- Slide 8 (index 7): Workflow Architecture ---
    # Mostly an image — just ensure title is correct
    slide8 = prs.slides[7]
    for shape in iter_all_shapes(slide8):
        if not shape.has_text_frame:
            continue
        replace_in_runs(shape, "ComplianceGPT", "SentriX")

    # --- Slide 10 (index 9): AI Agents ---
    slide10 = prs.slides[9]
    for shape in iter_all_shapes(slide10):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if "Clause Mapping Agent consists" in t:
            set_multiformat_text(shape, [
                ("SentriX's 7-Agent Pipeline orchestrated via Agent Builder", GREEN, True),
            ])

    # --- Slide 11 (index 10): Copilot features ---
    slide11 = prs.slides[10]
    for shape in iter_all_shapes(slide11):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        if "Natural Language Interface" in t:
            set_multiformat_text(shape, [
                ("Natural Language Interface", GREEN, True),
                (": No technical expertise needed — ask compliance questions in plain English and get instant, actionable answers.", WHITE, True),
            ])

        if "Assessment-Aware" in t:
            set_multiformat_text(shape, [
                ("Assessment-Aware Intelligence", GREEN, True),
                (": The Copilot contextualizes every answer with your organization's specific assessment data, scores, and gaps.", WHITE, True),
            ])

        if "One-click access" in t:
            set_multiformat_text(shape, [
                ("One-Click Deep Insights", GREEN, True),
                (": Score Drivers (why scores are what they are), Executive Summary (board-ready overview), Critical Clauses (highest risk), Remediation Plan (prioritized action items).", WHITE, True),
            ])

    # --- Slide 12 (index 11): Impact ---
    slide12 = prs.slides[11]
    for shape in iter_all_shapes(slide12):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        if t == "IMPACT":
            pass  # Keep as-is

        if "Zero Evidence Blind Spots" in t:
            set_multiformat_text(shape, [
                ("Zero Evidence Blind Spots", GREEN, True),
            ])

        if "Bias-Free" in t and "Scoring" in t and len(t) < 30:
            set_multiformat_text(shape, [
                ("Bias-Free Reproducible Scoring", GREEN, True),
            ])

        if "Before & After" in t:
            set_multiformat_text(shape, [
                ("The Problem", AMBER, True),
                (": ", WHITE, True),
                ("Two human auditors + same document = two completely different compliance scores. ", WHITE, False),
                ("Our 3-Tier Hybrid Scoring Engine", GREEN, True),
                (" (ML Semantic + Groq AI + Keyword NLP) produces ", WHITE, False),
                ("mathematically identical results every time", GREEN, False),
                (". Every score is traceable to specific clause evidence — fully audit-ready.", WHITE, False),
            ])

        if "Reality:" in t and "You can't confidently" in t:
            set_multiformat_text(shape, [
                ("\n", WHITE, False),
                ("The Truth", AMBER, True),
                (": ", WHITE, True),
                ("You can't fix what you can't measure accurately. Every score SentriX produces is ", WHITE, False),
                ("mathematically reproducible", GREEN, True),
                (", strictly tied to ISO clauses, and ", WHITE, False),
                ("completely audit-proof", GREEN, True),
                (". Our Evidence Validation Agent catches exactly what human auditors look for — ", WHITE, False),
                ("weak evidence, vague references, and missing documentation", GREEN, False),
                (" — before the real audit.", WHITE, False),
            ])

        if "Compliance Officers:" in t:
            set_multiformat_text(shape, [
                ("Compliance Officers: ", AMBER, True),
                ("SentriX handles 4-8 weeks of document analysis in minutes, freeing them to focus on actual risk strategy and remediation execution.", WHITE, False),
            ])

        if "Risk Managers:" in t:
            set_multiformat_text(shape, [
                ("Risk Managers: ", AMBER, True),
                ("Move from subjective 'gut feelings' to mathematical proof. The 3-Tier engine provides consistent, reproducible scores defensible to any external auditor.", WHITE, False),
            ])

        if "C-Suite Executives:" in t:
            set_multiformat_text(shape, [
                ("C-Suite Executives: ", AMBER, True),
                ("Instant clarity via a live dashboard showing the company's exact compliance health across all 6 standards simultaneously — no more waiting months for a static PDF.", WHITE, False),
            ])

        if "Scaling Organizations:" in t:
            set_multiformat_text(shape, [
                ("Scaling Organizations: ", AMBER, True),
                ("Cross-Standard Synergy Detection identifies where one policy fix satisfies multiple ISO standards at once, eliminating duplicated effort and cost.", WHITE, False),
            ])

    # --- Slide 13 (index 12): Strategic Plan ---
    slide13 = prs.slides[12]
    for shape in iter_all_shapes(slide13):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if "Strategic Development" in t:
            set_multiformat_text(shape, [
                ("Strategic Development Roadmap & Implementation Plan", GREEN, True),
            ])

    # --- Slide 14 (index 13): Final User Sees ---
    slide14 = prs.slides[13]
    for shape in iter_all_shapes(slide14):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()
        if "What the Final User Sees" in t:
            set_multiformat_text(shape, [
                ("What the User Experiences: SentriX in Action", GREEN, True),
            ])

    # --- Slide 15 (index 14): Code Links ---
    slide15 = prs.slides[14]
    for shape in iter_all_shapes(slide15):
        if not shape.has_text_frame:
            continue
        t = shape.text.strip()

        if "</> Code Links </>" in t:
            set_multiformat_text(shape, [
                ("  ", GREEN, True),
                ("</>", DGREEN, True),
                (" Code Repository ", GREEN, True),
                ("</>", DGREEN, True),
            ])

        if "Home Page Code" in t:
            set_multiformat_text(shape, [
                ("Full Platform Code :—      ", GREEN, True),
                ("https://github.com/MyGitHubProfile/SentriX-platform", BLUE, True),
            ])

        if "Agents Code" in t and shape.name == "TextBox 8":
            set_multiformat_text(shape, [
                ("AI Agents Code :—", GREEN, True),
            ])

        if "Agents Code" in t and shape.name == "TextBox 9":
            set_multiformat_text(shape, [
                ("Documentation :—", GREEN, True),
            ])

        if "github.com" in t and "agents" in t:
            set_multiformat_text(shape, [
                ("     ", GREEN, True),
                ("https://github.com/MyGitHubProfile/SentriX-platform/tree/main/server/src/agents", BLUE, True),
            ])

    # ===== SAVE =====
    prs.save(OUTPUT)
    print(f"Presentation saved to: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
